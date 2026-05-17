"""
Production-grade PPTX slide generator for analysis reports.
Applies DesignIntelligence agent output to create rich, professional presentations.

Architecture:
  - ThemeEngine: color tokens, typography scales, spacing grid
  - ShapeKit: decorative primitives (gradients, patterns, dividers, icon-shapes)
  - SlideBuilder: templated slide constructors
  - generate_slides(): top-level orchestrator
"""

from __future__ import annotations

import os
import math
import re
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

from app.utils.logger import get_logger

logger = get_logger(__name__)

# ── Canvas ───────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
GRID = Inches(0.4)       # base spacing unit
MARGIN = Inches(0.8)
CONTENT_W = SLIDE_W - MARGIN * 2

# ── Pre-built palettes (extended from design_intelligence) ──
THEME_BG: Dict[str, str] = {
    "finance": "0A1F39", "healthcare": "0D2B1D", "sales": "3D1A00",
    "generic": "1A1A2E", "tech": "0D1120", "executive": "141414",
}

LIGHT_BG: Dict[str, str] = {
    "finance": "F6FAFF", "healthcare": "F4FBF7", "sales": "FFF6F0",
    "generic": "F7F8FB", "tech": "F4F6FF", "executive": "FAFAFA",
}

STYLE_PRESETS: Dict[str, Dict[str, str]] = {
    "aurora": {"mode": "dark"},
    "editorial": {"mode": "light"},
    "minimal": {"mode": "light"},
    "bold": {"mode": "dark"},
}

THEME_STYLE_FALLBACK: Dict[str, str] = {
    "finance": "editorial",
    "healthcare": "minimal",
    "sales": "bold",
    "tech": "aurora",
    "executive": "minimal",
    "generic": "editorial",
}


# ═══════════════════════════════════════════════════════════
# Theme Engine
# ═══════════════════════════════════════════════════════════

class ThemeEngine:
    """Color tokens, typography, and spacing for a presentation."""

    def __init__(self, palette: List[str], theme_name: str, mode: str = "dark"):
        if len(palette) < 4:
            palette = ["#4f81bd", "#9cbb58", "#f79646", "#8064a2"]
        self.p = [_hex(palette[i]) for i in range(4)]
        bg_hex = THEME_BG.get(theme_name, "1A1A2E") if mode == "dark" else LIGHT_BG.get(theme_name, "F7F8FB")
        self.bg = _hex(bg_hex)

        if mode == "dark":
            self.white = _hex("FFFFFF")
            self.off_white = _hex("E8ECF0")
            self.muted = _hex("8899AA")
            self.dark_muted = _hex("556677")
            surface_alpha = 0.12
            surface_light_alpha = 0.08
        else:
            self.white = _hex("1E2430")
            self.off_white = _hex("2B3440")
            self.muted = _hex("5C6B78")
            self.dark_muted = _hex("7B8894")
            surface_alpha = 0.08
            surface_light_alpha = 0.05

        self.surface = _blend_on_bg(self.p[0], self.bg, surface_alpha)
        self.surface_light = _blend_on_bg(self.p[1], self.bg, surface_light_alpha)

    def accent(self, i: int = 0) -> RGBColor:
        return self.p[i % 4]

    def gradient_stops(self, c1: RGBColor, c2: RGBColor) -> List[RGBColor]:
        """Return 6 interpolated stops between two colors."""
        return [
            RGBColor(
                int(c1[0] + (c2[0] - c1[0]) * t),
                int(c1[1] + (c2[1] - c1[1]) * t),
                int(c1[2] + (c2[2] - c1[2]) * t),
            )
            for t in [0.0, 0.2, 0.4, 0.6, 0.8, 1.0]
        ]


# ═══════════════════════════════════════════════════════════
# Shape Kit — decorative primitives
# ═══════════════════════════════════════════════════════════

def _add_rect(slide, x, y, w, h, fill: RGBColor, alpha: float = 1.0, bg: RGBColor | None = None) -> Any:
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    if alpha >= 1.0:
        s.fill.fore_color.rgb = fill
    elif bg is not None:
        s.fill.fore_color.rgb = _blend_on_bg(fill, bg, alpha)
    else:
        s.fill.fore_color.rgb = _blend_to_dark(fill, alpha)
    s.line.fill.background()
    return s


def _add_oval(slide, x, y, w, h, fill: RGBColor, alpha: float = 1.0, bg: RGBColor | None = None) -> Any:
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid()
    if alpha >= 1.0:
        s.fill.fore_color.rgb = fill
    elif bg is not None:
        s.fill.fore_color.rgb = _blend_on_bg(fill, bg, alpha)
    else:
        s.fill.fore_color.rgb = _blend_to_dark(fill, alpha)
    s.line.fill.background()
    return s


def _blend_on_bg(color: RGBColor, bg: RGBColor, alpha: float) -> RGBColor:
    """Blend a color onto a background color to approximate transparency."""
    return RGBColor(
        max(0, min(255, int(bg[0] + (color[0] - bg[0]) * alpha))),
        max(0, min(255, int(bg[1] + (color[1] - bg[1]) * alpha))),
        max(0, min(255, int(bg[2] + (color[2] - bg[2]) * alpha))),
    )


def _blend_to_dark(color: RGBColor, alpha: float) -> RGBColor:
    """Approximate a transparent color on dark background by blending to near-black."""
    return RGBColor(
        max(0, min(255, int(color[0] * alpha))),
        max(0, min(255, int(color[1] * alpha))),
        max(0, min(255, int(color[2] * alpha))),
    )


def _add_text_box(slide, x, y, w, h, text: str, size: int, color: RGBColor,
                  bold: bool = False, align=PP_ALIGN.LEFT, font: str = "Calibri",
                  italic: bool = False, spacing: int = 0) -> Any:
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    if spacing:
        p.line_spacing = Pt(spacing)
    return tb


def _add_multi_text(slide, x, y, w, h, lines: List[Tuple[str, int, RGBColor, bool]],
                    font: str = "Calibri", align=PP_ALIGN.LEFT, spacing: int = 26) -> Any:
    """Add text box with multiple styled paragraphs."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        p.line_spacing = Pt(spacing)
    return tb


# ── Decorative elements ──────────────────────────────────

def _geometric_background(slide, theme: ThemeEngine):
    """Layer abstract geometric shapes for visual depth on dark slides."""
    # Large faint circle top-right
    _add_oval(slide, Inches(10.5), Inches(-1.5), Inches(5), Inches(5),
              theme.accent(0), alpha=0.04, bg=theme.bg)
    # Small accent circle bottom-left
    _add_oval(slide, Inches(-0.5), Inches(6.0), Inches(2.5), Inches(2.5),
              theme.accent(1), alpha=0.06, bg=theme.bg)
    # Horizontal accent line near top
    _add_rect(slide, MARGIN, Inches(0.2), Inches(2), Inches(0.015),
              theme.accent(0), alpha=0.8, bg=theme.bg)
    # Subtle grid dot bottom-right
    _add_oval(slide, Inches(12.2), Inches(6.8), Inches(0.25), Inches(0.25),
              theme.accent(2), alpha=0.15, bg=theme.bg)


def _light_geometric_background(slide, theme: ThemeEngine):
    """Subtle geometric accents for light slides."""
    # Thin accent bar top
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04), theme.accent(0))
    # Small corner accent
    _add_rect(slide, Inches(0), Inches(0), Inches(0.4), Inches(0.04), theme.accent(1))
    # Faint dot bottom-right
    _add_oval(slide, Inches(12.5), Inches(7.0), Inches(0.3), Inches(0.3),
              theme.accent(0), alpha=0.06, bg=theme.bg)


def _apply_motif(slide, theme: ThemeEngine, motif: str):
    """Apply a theme motif with subtle background accents."""
    motif_key = (motif or "").lower()
    if motif_key == "grid":
        _motif_grid(slide, theme)
    elif motif_key == "rings":
        _motif_rings(slide, theme)
    elif motif_key == "diagonal":
        _motif_diagonal(slide, theme)
    elif motif_key == "dots":
        _motif_dots(slide, theme)
    elif motif_key == "line":
        _motif_line(slide, theme)


def _motif_grid(slide, theme: ThemeEngine):
    for i in range(1, 6):
        y = Inches(0.8 + i * 1.0)
        _add_rect(slide, MARGIN, y, SLIDE_W - MARGIN * 2, Inches(0.01),
                  theme.accent(0), alpha=0.06, bg=theme.bg)
    for i in range(1, 6):
        x = MARGIN + Inches(1.8) * i
        _add_rect(slide, x, Inches(0.8), Inches(0.01), SLIDE_H - Inches(1.6),
                  theme.accent(1), alpha=0.04, bg=theme.bg)


def _motif_rings(slide, theme: ThemeEngine):
    _add_oval(slide, Inches(9.2), Inches(-1.8), Inches(5.5), Inches(5.5),
              theme.accent(0), alpha=0.04, bg=theme.bg)
    _add_oval(slide, Inches(8.6), Inches(-1.2), Inches(4.5), Inches(4.5),
              theme.accent(1), alpha=0.05, bg=theme.bg)


def _motif_diagonal(slide, theme: ThemeEngine):
    bar = _add_rect(slide, Inches(-1.0), Inches(5.5), Inches(7.5), Inches(0.35),
                    theme.accent(0), alpha=0.12, bg=theme.bg)
    bar.rotation = -12
    bar2 = _add_rect(slide, Inches(6.0), Inches(1.2), Inches(8.0), Inches(0.25),
                     theme.accent(1), alpha=0.08, bg=theme.bg)
    bar2.rotation = 10


def _motif_dots(slide, theme: ThemeEngine):
    for i in range(6):
        _add_oval(slide, Inches(10.8 + i * 0.35), Inches(1.0 + i * 0.25), Inches(0.12), Inches(0.12),
                  theme.accent(i), alpha=0.18, bg=theme.bg)


def _motif_line(slide, theme: ThemeEngine):
    _add_rect(slide, MARGIN, Inches(6.6), Inches(3.5), Inches(0.03),
              theme.accent(0), alpha=0.4, bg=theme.bg)


def _resolve_template(design_spec: Dict[str, Any], theme_name: str) -> Dict[str, str]:
    style = (design_spec.get("template_style") or "").lower()
    if style not in STYLE_PRESETS:
        style = THEME_STYLE_FALLBACK.get(theme_name, "editorial")
    profile = dict(STYLE_PRESETS.get(style, {"mode": "dark"}))
    profile["name"] = style
    return profile


def _default_typography(style_name: str) -> Dict[str, str]:
    if style_name == "editorial":
        return {"title_font": "Georgia", "body_font": "Calibri", "stat_font": "Arial Black"}
    if style_name == "minimal":
        return {"title_font": "Calibri Light", "body_font": "Calibri", "stat_font": "Calibri"}
    if style_name == "bold":
        return {"title_font": "Arial Black", "body_font": "Arial", "stat_font": "Arial Black"}
    return {"title_font": "Trebuchet MS", "body_font": "Calibri", "stat_font": "Arial Black"}


def _resolve_typography(design_spec: Dict[str, Any], style_name: str) -> Dict[str, str]:
    typography = design_spec.get("typography") or {}
    defaults = _default_typography(style_name)
    return {
        "title_font": typography.get("title_font") or defaults["title_font"],
        "body_font": typography.get("body_font") or defaults["body_font"],
        "stat_font": typography.get("stat_font") or defaults["stat_font"],
    }


def _accent_divider(slide, y, theme: ThemeEngine, width_ratio: float = 0.15):
    """Thin accent line at a given y position."""
    w = Inches(2.5 * width_ratio / 0.15)
    _add_rect(slide, MARGIN, y, w, Inches(0.025), theme.accent(0))


def _stat_card(slide, x, y, w, h, number: str, label: str, theme: ThemeEngine, idx: int):
    """Rendered stat callout card with number + label."""
    color = theme.accent(idx % 4)
    # Card background
    _add_rect(slide, x, y, w, h, theme.surface)
    # Top accent bar
    _add_rect(slide, x, y, w, Inches(0.05), color)
    # Number
    _add_text_box(slide, x + Inches(0.3), y + Inches(0.25), w - Inches(0.6), Inches(0.7),
                  number, 48, color, bold=True, align=PP_ALIGN.LEFT)
    # Label
    _add_text_box(slide, x + Inches(0.3), y + h - Inches(0.55), w - Inches(0.6), Inches(0.4),
                  label, 12, theme.muted, align=PP_ALIGN.LEFT)


def _icon_shape(slide, x, y, size, theme: ThemeEngine, idx: int, shape_type=MSO_SHAPE.OVAL):
    """Decorative icon placeholder using geometric shapes."""
    s = slide.shapes.add_shape(shape_type, x, y, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = _blend_on_bg(theme.accent(idx), theme.bg, 0.2)
    s.line.fill.background()
    # Inner shape
    inner_size = size * 0.5
    inner_x = x + (size - inner_size) // 2
    inner_y = y + (size - inner_size) // 2
    s2 = slide.shapes.add_shape(shape_type, inner_x, inner_y, inner_size, inner_size)
    s2.fill.solid()
    s2.fill.fore_color.rgb = _blend_on_bg(theme.accent(idx), theme.bg, 0.15)
    s2.line.fill.background()


def _side_accent_bar(slide, theme: ThemeEngine):
    """Vertical accent bar on the left side of the slide."""
    _add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, theme.accent(0))


def _bottom_strip(slide, theme: ThemeEngine, alpha: float = 0.85):
    """Subtle colored strip at the bottom."""
    _add_rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), theme.accent(0), alpha=alpha, bg=theme.bg)
    # Small color blocks
    for i in range(4):
        _add_rect(slide, Inches(0.8 + i * 0.4), Inches(7.3), Inches(0.28), Inches(0.06),
                  theme.accent(i))


# ═══════════════════════════════════════════════════════════
# Slide Builders
# ═══════════════════════════════════════════════════════════

def _title_slide(prs, title: str, subtitle: str, theme: ThemeEngine, palette: List[str], brief: dict,
                 style: str, typography: Dict[str, str]):
    """Hero / cover slide with big typography and geometric depth."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, theme.bg)

    title_font = typography["title_font"]
    body_font = typography["body_font"]
    stat_font = typography["stat_font"]
    mood = brief.get("mood_description", "")[:120] or "AI-Powered Insights & Recommendations"
    subtitle = subtitle[:150] or mood
    motif = brief.get("visual_motif", "")

    if style in ("editorial", "minimal"):
        _light_geometric_background(slide, theme)
        _apply_motif(slide, theme, motif)
        _side_accent_bar(slide, theme)

        _add_text_box(slide, MARGIN + Inches(0.3), Inches(1.0), Inches(6), Inches(0.4),
                      "DATA ANALYSIS REPORT", 10, theme.accent(0), bold=True, font=body_font)
        _add_text_box(slide, MARGIN + Inches(0.3), Inches(1.6), Inches(10.5), Inches(1.75),
                      title[:110], 42, theme.white, bold=True, font=title_font, spacing=52)
        _add_text_box(slide, MARGIN + Inches(0.3), Inches(3.72), Inches(8.5), Inches(0.55),
                      subtitle, 15, theme.off_white, font=body_font, spacing=20)
        _add_text_box(slide, MARGIN + Inches(0.3), Inches(4.32), Inches(8.5), Inches(0.35),
                      mood, 10, theme.muted, italic=True, font=body_font)

        _accent_divider(slide, Inches(4.72), theme, width_ratio=0.2)
        _add_text_box(slide, MARGIN + Inches(0.3), Inches(5.02), Inches(6), Inches(0.3),
                      f"Theme: {brief.get('theme', 'Professional')}  •  Template: {style}",
                      9, theme.dark_muted, font=body_font)
        return

    if style == "bold":
        _geometric_background(slide, theme)
        _apply_motif(slide, theme, motif)
        _add_rect(slide, MARGIN, Inches(4.4), SLIDE_W - MARGIN * 2, Inches(2.4),
                  theme.accent(0), alpha=0.14, bg=theme.bg)

        _add_text_box(slide, MARGIN, Inches(1.1), Inches(6), Inches(0.4),
                      "DATA ANALYSIS REPORT", 11, theme.accent(0), bold=True, font=body_font)
        _add_text_box(slide, MARGIN, Inches(1.7), Inches(11.0), Inches(2.0),
                      title[:100], 46, theme.white, bold=True, font=stat_font, spacing=56)
        _add_text_box(slide, MARGIN, Inches(3.95), Inches(8.8), Inches(0.55),
                      subtitle, 16, theme.off_white, font=body_font, spacing=22)
        _add_text_box(slide, MARGIN, Inches(4.6), Inches(8), Inches(0.35),
                      mood, 10, theme.muted, italic=True, font=body_font)
        _bottom_strip(slide, theme, alpha=0.45 if style in ("editorial", "minimal") else 0.85)
        return

    _geometric_background(slide, theme)
    _apply_motif(slide, theme, motif)

    # Large decorative circle behind title area
    _add_oval(slide, Inches(-3), Inches(-2), Inches(8), Inches(8), theme.accent(0), alpha=0.04, bg=theme.bg)
    _add_oval(slide, Inches(9), Inches(2), Inches(7), Inches(7), theme.accent(1), alpha=0.03, bg=theme.bg)

    # Top accent line
    _add_rect(slide, MARGIN, Inches(0.8), Inches(3.5), Inches(0.04), theme.accent(0))

    # Category label
    _add_text_box(slide, MARGIN, Inches(1.1), Inches(6), Inches(0.4),
                  "DATA ANALYSIS REPORT", 11, theme.accent(0), bold=True,
                  font=body_font, spacing=0)

    # Main title
    clean_title = title[:100]
    _add_text_box(slide, MARGIN, Inches(1.7), Inches(11.5), Inches(2.0),
                  clean_title, 44, theme.white, bold=True, font=title_font, spacing=56)

    # Subtitle line
    _add_text_box(slide, MARGIN, Inches(3.85), Inches(8.8), Inches(0.55),
                  subtitle, 16, theme.off_white, font=body_font, spacing=22)
    _add_text_box(slide, MARGIN, Inches(4.45), Inches(8), Inches(0.35),
                  mood, 10, theme.muted, italic=True, font=body_font)

    # Divider
    _accent_divider(slide, Inches(5.0), theme)

    # Bottom metadata
    _add_text_box(slide, MARGIN, Inches(5.3), Inches(5), Inches(0.3),
                  f"Theme: {brief.get('theme', 'Professional')}  •  Template: {style}",
                  10, theme.dark_muted, font=body_font)

    _bottom_strip(slide, theme, alpha=0.45 if style in ("editorial", "minimal") else 0.85)


def _executive_summary_slide(prs, narrative: str, theme: ThemeEngine, sections: list,
                             style: str, typography: Dict[str, str]):
    """Executive summary with key points extracted from narrative."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, theme.bg)

    title_font = typography["title_font"]
    body_font = typography["body_font"]

    if style in ("editorial", "minimal"):
        _light_geometric_background(slide, theme)
        _side_accent_bar(slide, theme)

        _section_label(slide, "EXECUTIVE SUMMARY", theme, body_font)
        _add_text_box(slide, MARGIN, Inches(0.65), Inches(8.5), Inches(0.7),
                      "Key Findings at a Glance", 26, theme.white, bold=True, font=title_font)
        _accent_divider(slide, Inches(1.4), theme, width_ratio=0.18)

        para = narrative.strip()[:900] if narrative else "Analysis completed. See insights below."
        _add_text_box(slide, MARGIN, Inches(1.8), Inches(7.8), Inches(1.2),
                      para, 14, theme.off_white, font=body_font, spacing=24)

        if sections:
            _add_text_box(slide, MARGIN, Inches(3.4), Inches(4), Inches(0.4),
                          "REPORT SECTIONS", 10, theme.accent(0), bold=True, font=body_font)
            sec_list = "\n".join(
                f"▸  {s.get('title', 'Section')}" for s in sections[:4]
            )
            _add_text_box(slide, MARGIN, Inches(3.8), Inches(6), Inches(2.5),
                          sec_list, 12, theme.muted, font=body_font, spacing=22)

        _add_text_box(slide, Inches(8.2), Inches(3.2), Inches(4), Inches(0.4),
                      "DESIGN PALETTE", 10, theme.accent(0), bold=True, font=body_font)
        for i, c in enumerate(theme.p[:4]):
            _add_rect(slide, Inches(8.2 + i * 0.7), Inches(3.7), Inches(0.5), Inches(0.5), c)
            _add_text_box(slide, Inches(8.2 + i * 0.7), Inches(4.25), Inches(0.5), Inches(0.2),
                          f"#{_to_hex(c)}", 7, theme.dark_muted, align=PP_ALIGN.CENTER, font=body_font)
        return

    _side_accent_bar(slide, theme)

    _section_label(slide, "EXECUTIVE SUMMARY", theme, body_font)
    _add_text_box(slide, MARGIN, Inches(0.65), Inches(8), Inches(0.7),
                  "Key Findings at a Glance", 28, theme.white, bold=True, font=title_font)

    _accent_divider(slide, Inches(1.5), theme)

    para = narrative.strip()[:800] if narrative else "Analysis completed. See insights below."
    _add_text_box(slide, MARGIN, Inches(1.8), Inches(7.5), Inches(1.2),
                  para, 15, theme.off_white, font=body_font, spacing=24)

    if sections:
        _add_text_box(slide, MARGIN, Inches(3.3), Inches(4), Inches(0.4),
                      "REPORT SECTIONS", 10, theme.accent(0), bold=True, font=body_font)
        sec_list = "\n".join(
            f"▸  {s.get('title', 'Section')}" for s in sections[:4]
        )
        _add_text_box(slide, MARGIN, Inches(3.7), Inches(6), Inches(2.5),
                      sec_list, 13, theme.muted, font=body_font, spacing=22)

    _add_text_box(slide, Inches(8.5), Inches(3.3), Inches(4), Inches(0.4),
                  "DESIGN PALETTE", 10, theme.accent(0), bold=True, font=body_font)
    for i, c in enumerate(theme.p[:4]):
        _add_rect(slide, Inches(8.5 + i * 0.7), Inches(3.8), Inches(0.5), Inches(0.5), c)
        _add_text_box(slide, Inches(8.5 + i * 0.7), Inches(4.35), Inches(0.5), Inches(0.2),
                      f"#{_to_hex(c)}", 7, theme.dark_muted, align=PP_ALIGN.CENTER, font=body_font)

    _bottom_strip(slide, theme, alpha=0.45 if style in ("editorial", "minimal") else 0.85)


def _story_arc_slide(prs, design_spec: dict, theme: ThemeEngine, style: str, typography: Dict[str, str]):
    """Visual storytelling roadmap — shows the narrative journey before diving into insights."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, theme.bg)

    title_font = typography["title_font"]
    body_font = typography["body_font"]
    arc = design_spec.get("storytelling_arc", "Context → Discovery → Evidence → Action")
    principle = design_spec.get("design_principle", "")

    motif = design_spec.get("visual_motif", "")

    if style in ("editorial", "minimal"):
        _light_geometric_background(slide, theme)
    _apply_motif(slide, theme, motif)
    _side_accent_bar(slide, theme)

    _section_label(slide, "NARRATIVE JOURNEY", theme, body_font)
    _add_text_box(slide, MARGIN, Inches(0.65), Inches(10), Inches(0.7),
                  "The Data Story", 28, theme.white, bold=True, font=title_font)
    _add_text_box(slide, MARGIN, Inches(1.48), Inches(9), Inches(0.4),
                  principle[:150] if principle else "Building insight through systematic data exploration.",
                  12, theme.muted, italic=True, font=body_font)

    _accent_divider(slide, Inches(2.08), theme)

    # Arc phases as horizontal milestone cards
    phases = [p.strip() for p in arc.split("→")]
    if not phases or phases == [arc]:
        phases = ["Context", "Discovery", "Evidence", "Action"]

    card_w = Inches(2.6)
    card_h = Inches(3.2)
    gap = Inches(0.4)
    total_w = len(phases) * card_w + (len(phases) - 1) * gap
    start_x = (SLIDE_W - total_w) // 2
    card_y = Inches(2.65)

    phase_icons = ["🔍", "📊", "📈", "🎯", "💡", "🚀"]
    phase_labels = [
        "UNDERSTAND\nWhat the data contains", "DISCOVER\nPatterns & relationships",
        "ANALYZE\nDeep statistical tests", "VALIDATE\nSignificance & confidence",
        "SYNTHESIZE\nKey findings & narrative", "ACT\nRecommendations & next steps"
    ]

    for i, phase in enumerate(phases[:6]):
        cx = start_x + i * (card_w + gap)
        # Card background
        _add_rect(slide, cx, card_y, card_w, card_h, theme.surface)
        # Top accent
        _add_rect(slide, cx, card_y, card_w, Inches(0.05), theme.accent(i))
        # Phase number
        _add_text_box(slide, cx + Inches(0.2), card_y + Inches(0.2), Inches(1), Inches(0.4),
                  f"0{i + 1}", 14, theme.accent(i), bold=True, font=body_font)
        # Phase name
        _add_text_box(slide, cx + Inches(0.2), card_y + Inches(0.7), card_w - Inches(0.4), Inches(0.6),
                  phase, 14, theme.white, bold=True, font=title_font)
        # Phase description
        label_text = phase_labels[i] if i < len(phase_labels) else phase
        _add_text_box(slide, cx + Inches(0.2), card_y + Inches(1.5), card_w - Inches(0.4), Inches(1.5),
                  label_text, 10, theme.muted, font=body_font, spacing=16)
        # Connector arrow (except last)
        if i < len(phases) - 1:
            arrow_x = cx + card_w
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x, card_y + card_h // 2 - Inches(0.12),
                gap, Inches(0.24),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = _blend_on_bg(theme.accent(i), theme.bg, 0.3)
            arrow.line.fill.background()

    _bottom_strip(slide, theme, alpha=0.45 if style in ("editorial", "minimal") else 0.85)


def _insight_slides(prs, insights: list, chart_paths: list, theme: ThemeEngine, palette: list,
                    design_spec: dict, style: str, typography: Dict[str, str],
                    insight_limit: int | None = None):
    """One insight per slide — with storytelling context and visual chart integration."""
    arc = design_spec.get("storytelling_arc", "Context → Discovery → Evidence → Action")
    title_font = typography["title_font"]
    body_font = typography["body_font"]
    stat_font = typography["stat_font"]
    motif = design_spec.get("visual_motif", "")

    selected_insights = insights[:insight_limit] if insight_limit else insights
    for idx, ins in enumerate(selected_insights):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, theme.bg)
        if style in ("editorial", "minimal"):
            _light_geometric_background(slide, theme)
        else:
            _geometric_background(slide, theme)
        _apply_motif(slide, theme, motif)

        # Compact slide number marker. Keep it away from the main copy.
        num = f"{idx + 1:02d}"
        txb = slide.shapes.add_textbox(Inches(11.05), Inches(0.35), Inches(1.4), Inches(0.7))
        tf = txb.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = _blend_on_bg(theme.accent(idx % 4), theme.bg, 0.35)
        p.font.name = stat_font
        p.alignment = PP_ALIGN.RIGHT

        # Section label with storytelling phase
        _section_label(slide, f"INSIGHT {idx + 1}  •  {arc}", theme, body_font)

        # LEFT COLUMN: Text content
        content = _clean_markdown_text(ins.get("content", ""))
        score = ins.get("significance_score", 0)
        sources = ins.get("source_agents", [])

        # Main insight statement
        _add_text_box(slide, MARGIN, Inches(1.2), Inches(6.8), Inches(1.8),
                  content, 22, theme.white, bold=True, font=title_font, spacing=34)

        _accent_divider(slide, Inches(3.0), theme, width_ratio=0.12)

        # Storytelling context
        story_context = _story_context_for_insight(idx, content, arc)
        _add_text_box(slide, MARGIN, Inches(3.3), Inches(6.8), Inches(0.8),
                  story_context, 12, theme.muted, italic=True, font=body_font, spacing=18)

        # Significance score bar
        bar_y = Inches(4.3)
        _add_rect(slide, MARGIN, bar_y, Inches(3.5), Inches(0.05), theme.dark_muted)
        _add_rect(slide, MARGIN, bar_y, Inches(3.5 * score), Inches(0.05), theme.accent(idx % 4))
        _add_text_box(slide, Inches(4.5), Inches(4.2), Inches(1.2), Inches(0.3),
                  f"{score:.0%}", 12, theme.accent(idx % 4), bold=True, font=stat_font)
        _add_text_box(slide, MARGIN, Inches(4.5), Inches(4), Inches(0.3),
                  f"Confidence  •  {', '.join(sources[:2]) if sources else 'Multiple agents'}",
                  9, theme.dark_muted, font=body_font)

        # Evidence data points
        evidence = ins.get("data_evidence", {})
        if evidence:
            ev_items = []
            for k, v in list(evidence.items())[:4]:
                ev_items.append(f"▸  {k}: {v}")
            _add_text_box(slide, MARGIN, Inches(5.0), Inches(6.8), Inches(2.0),
                          "\n".join(ev_items), 10, theme.off_white, font=body_font, spacing=16)

        # RIGHT COLUMN: Visual chart reference
        right_x = Inches(7.8)
        # Chart image area
        chart_img = _find_chart_for_insight(idx, chart_paths)
        if chart_img:
            try:
                _add_rect(slide, right_x - Inches(0.08), Inches(1.12), Inches(4.96), Inches(3.96), theme.surface)
                slide.shapes.add_picture(chart_img, right_x, Inches(1.2), Inches(4.8), Inches(3.8))
            except Exception:
                _chart_placeholder(slide, right_x, Inches(1.2), Inches(4.8), Inches(3.8),
                                   theme, idx, "Chart visualization", body_font)
        else:
            # Decorative placeholder with insight number
            _chart_placeholder(slide, right_x, Inches(1.2), Inches(4.8), Inches(3.8),
                               theme, idx, f"Insight {idx + 1}", body_font)

        # Right-side insight tag
        tag_y = Inches(5.3)
        _add_rect(slide, right_x, tag_y, Inches(4.8), Inches(0.04), theme.accent(idx % 4))
        _add_text_box(slide, right_x, tag_y + Inches(0.15), Inches(4.8), Inches(0.5),
                  f"Key Insight #{idx + 1}: {content[:100]}...",
                  9, theme.muted, font=body_font, spacing=14)

        _bottom_strip(slide, theme, alpha=0.45 if style in ("editorial", "minimal") else 0.85)


def _story_context_for_insight(idx: int, content: str, arc: str) -> str:
    """Generate a short storytelling bridge for an insight based on its position in the arc."""
    contexts = [
        "Setting the foundation — this finding establishes the baseline understanding of the data landscape.",
        "Drilling deeper — uncovering the structural patterns that drive the observed outcomes.",
        "Connecting the dots — revealing how different variables interact and influence each other.",
        "Statistical validation — confirming whether these patterns are signal or noise.",
        "Quantifying impact — measuring the magnitude and practical significance of key effects.",
        "Synthesizing across dimensions — this multi-faceted finding bridges multiple analytical perspectives.",
        "The actionable insight — translating statistical evidence into strategic direction.",
        "Looking forward — this finding points toward future trends and emerging opportunities.",
    ]
    return contexts[idx] if idx < len(contexts) else "A critical piece of the analytical mosaic."


def _find_chart_for_insight(idx: int, chart_paths: list) -> str | None:
    """Match an insight index to a corresponding chart image path."""
    png_paths = [_resolve_png_path(p) for p in chart_paths]
    valid = [p for p in png_paths if p and os.path.exists(p)]
    if idx < len(valid):
        return valid[idx]
    return None


def _chart_placeholder(slide, x, y, w, h, theme: ThemeEngine, idx: int, label: str, font: str):
    """Decorative chart area placeholder with geometric design."""
    # Border rectangle
    _add_rect(slide, x, y, w, h, theme.surface)
    # Inner decorative elements
    _icon_shape(slide, x + w // 2 - Inches(0.5), y + h // 2 - Inches(0.5), Inches(1.0), theme, idx)
    _add_text_box(slide, x + Inches(0.2), y + h - Inches(0.4), w - Inches(0.4), Inches(0.3),
                  label, 9, theme.dark_muted, align=PP_ALIGN.CENTER, font=font)


def _chart_slides(prs, chart_paths: list, theme: ThemeEngine, palette: list,
                  style: str, typography: Dict[str, str], motif: str = "",
                  chart_specs: list | None = None, insights: list | None = None,
                  chart_limit: int | None = None):
    """Chart deep-dive slides with editorial evidence layouts."""
    title_font = typography["title_font"]
    body_font = typography["body_font"]
    stat_font = typography["stat_font"]
    chart_specs = chart_specs or []
    insights = insights or []
    valid_paths = []
    for chart_path in chart_paths:
        png_path = _resolve_png_path(chart_path)
        if png_path and os.path.exists(png_path):
            valid_paths.append(png_path)

    for idx, png_path in enumerate(valid_paths[:chart_limit] if chart_limit else valid_paths):
        spec = chart_specs[idx] if idx < len(chart_specs) and isinstance(chart_specs[idx], dict) else {}
        related_insight = insights[idx] if idx < len(insights) and isinstance(insights[idx], dict) else {}
        narrative_role = _clean_markdown_text(spec.get("narrative_role") or _story_context_for_insight(idx, "", ""))
        takeaway = _clean_markdown_text(related_insight.get("content") or narrative_role)
        chart_name = _clean_markdown_text(spec.get("title") or os.path.splitext(os.path.basename(png_path))[0].replace("_", " ").title())
        layout_variant = idx % 3

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, theme.bg)
        if style in ("editorial", "minimal"):
            _light_geometric_background(slide, theme)
        else:
            _geometric_background(slide, theme)
        _apply_motif(slide, theme, motif)

        _section_label(slide, "DATA VISUALIZATION", theme, body_font)

        if layout_variant == 0:
            _add_text_box(slide, MARGIN, Inches(0.65), Inches(7.1), Inches(0.7),
                          chart_name[:92], 25, theme.white, bold=True, font=title_font)
            _accent_divider(slide, Inches(1.35), theme)
            _add_rect(slide, MARGIN - Inches(0.08), Inches(1.62), Inches(7.25), Inches(5.1), theme.surface)
            try:
                slide.shapes.add_picture(png_path, MARGIN, Inches(1.7), Inches(7.1), Inches(4.95))
            except Exception as exc:
                logger.log_error("Failed to embed chart image", exc, chart_path=png_path)
                _chart_placeholder(slide, MARGIN, Inches(1.7), Inches(7.1), Inches(4.95), theme, idx, chart_name, body_font)

            panel_x = Inches(8.35)
            _add_rect(slide, panel_x, Inches(1.45), Inches(4.0), Inches(5.25), theme.surface)
            _add_text_box(slide, panel_x + Inches(0.3), Inches(1.75), Inches(3.4), Inches(0.35),
                          "WHAT THIS PROVES", 9, theme.accent(idx), bold=True, font=body_font)
            _add_text_box(slide, panel_x + Inches(0.3), Inches(2.25), Inches(3.35), Inches(1.8),
                          takeaway[:230], 17, theme.white, bold=True, font=title_font, spacing=24)
            _add_text_box(slide, panel_x + Inches(0.3), Inches(4.35), Inches(3.35), Inches(1.2),
                          narrative_role[:210], 11, theme.muted, italic=True, font=body_font, spacing=17)
            _add_text_box(slide, panel_x + Inches(0.3), Inches(5.95), Inches(3.3), Inches(0.35),
                          f"Evidence view {idx + 1:02d}", 10, theme.dark_muted, font=stat_font)
        elif layout_variant == 1:
            _add_text_box(slide, MARGIN, Inches(0.65), Inches(11.7), Inches(0.55),
                          chart_name[:96], 23, theme.white, bold=True, font=title_font)
            _add_text_box(slide, MARGIN, Inches(1.38), Inches(10.5), Inches(0.45),
                          takeaway[:165], 12, theme.muted, italic=True, font=body_font)
            _add_rect(slide, MARGIN - Inches(0.08), Inches(2.02), Inches(11.66), Inches(4.35), theme.surface)
            try:
                slide.shapes.add_picture(png_path, MARGIN, Inches(2.1), Inches(11.5), Inches(4.2))
            except Exception as exc:
                logger.log_error("Failed to embed chart image", exc, chart_path=png_path)
                _chart_placeholder(slide, MARGIN, Inches(2.1), Inches(11.5), Inches(4.2), theme, idx, chart_name, body_font)
            _add_rect(slide, MARGIN, Inches(6.55), Inches(11.5), Inches(0.04), theme.accent(idx))
            _add_text_box(slide, MARGIN, Inches(6.67), Inches(9.8), Inches(0.35),
                          narrative_role[:180], 9, theme.dark_muted, font=body_font)
        else:
            _add_rect(slide, MARGIN, Inches(0.78), Inches(3.2), Inches(5.9), theme.surface)
            _add_text_box(slide, MARGIN + Inches(0.25), Inches(1.05), Inches(2.65), Inches(0.5),
                          f"{idx + 1:02d}", 30, theme.accent(idx), bold=True, font=stat_font)
            _add_text_box(slide, MARGIN + Inches(0.25), Inches(1.75), Inches(2.6), Inches(1.25),
                          chart_name[:70], 19, theme.white, bold=True, font=title_font, spacing=24)
            _add_text_box(slide, MARGIN + Inches(0.25), Inches(3.25), Inches(2.55), Inches(1.45),
                          narrative_role[:190], 11, theme.muted, font=body_font, spacing=17)
            _add_text_box(slide, MARGIN + Inches(0.25), Inches(5.45), Inches(2.55), Inches(0.65),
                          "Use this chart as the evidence anchor for the recommendation that follows.",
                          9, theme.dark_muted, italic=True, font=body_font, spacing=14)
            chart_x = Inches(4.35)
            _add_rect(slide, chart_x - Inches(0.08), Inches(0.95), Inches(8.1), Inches(5.75), theme.surface)
            try:
                slide.shapes.add_picture(png_path, chart_x, Inches(1.03), Inches(7.95), Inches(5.58))
            except Exception as exc:
                logger.log_error("Failed to embed chart image", exc, chart_path=png_path)
                _chart_placeholder(slide, chart_x, Inches(1.03), Inches(7.95), Inches(5.58), theme, idx, chart_name, body_font)

        _bottom_strip(slide, theme, alpha=0.45 if style in ("editorial", "minimal") else 0.85)


def _recommendations_slide(prs, sections: list, theme: ThemeEngine, palette: list,
                           style: str, typography: Dict[str, str], motif: str = ""):
    """Action-oriented recommendations with icon cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, theme.bg)

    title_font = typography["title_font"]
    body_font = typography["body_font"]

    if style in ("editorial", "minimal"):
        _light_geometric_background(slide, theme)
    else:
        _geometric_background(slide, theme)
    _apply_motif(slide, theme, motif)
    _side_accent_bar(slide, theme)

    _section_label(slide, "RECOMMENDATIONS", theme, body_font)
    _add_text_box(slide, MARGIN, Inches(0.65), Inches(8), Inches(0.7),
                  "Strategic Actions & Next Steps", 28, theme.white, bold=True, font=title_font)

    _accent_divider(slide, Inches(1.5), theme)

    # Try to find a recommendations section, otherwise use last section
    rec_section = next(
        (s for s in sections if "recommend" in s.get("title", "").lower()), None
    )
    content_text = ""
    if rec_section:
        content_text = rec_section.get("content", "")
    elif sections:
        content_text = sections[-1].get("content", "")

    # Parse bullet points from markdown content
    bullets = _extract_bullets(content_text)
    if not bullets:
        # Fallback: use insight content as bullets
        bullets = ["Review the full analysis for domain-specific actions",
                   "Share findings with key stakeholders",
                   "Set up periodic data reviews to track progress",
                   "Integrate insights into decision-making workflows"]

    # Display as action cards in a 2×2 grid
    card_w = Inches(5.5)
    card_h = Inches(2.2)
    positions = [
        (MARGIN, Inches(2.0)),
        (MARGIN + card_w + Inches(0.5), Inches(2.0)),
        (MARGIN, Inches(4.5)),
        (MARGIN + card_w + Inches(0.5), Inches(4.5)),
    ]

    for i, (bx, by) in enumerate(positions):
        if i >= len(bullets):
            break
        bullet = _clean_markdown_text(bullets[i])[:155]
        # Card
        _add_rect(slide, bx, by, card_w, card_h, theme.surface)
        # Left color bar
        _add_rect(slide, bx, by, Inches(0.08), card_h, theme.accent(i))
        # Number
        _add_text_box(slide, bx + Inches(0.35), by + Inches(0.2), Inches(0.5), Inches(0.5),
                  f"{i + 1}", 28, theme.accent(i), bold=True, font=title_font)
        # Bullet text
        _add_text_box(slide, bx + Inches(0.35), by + Inches(0.82), card_w - Inches(0.7), card_h - Inches(0.95),
                  bullet, 11, theme.off_white, font=body_font, spacing=16)

    _bottom_strip(slide, theme, alpha=0.45 if style in ("editorial", "minimal") else 0.85)


def _closing_slide(prs, theme: ThemeEngine, brief: dict, style: str, typography: Dict[str, str]):
    """Thank you / closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, theme.bg)

    title_font = typography["title_font"]
    body_font = typography["body_font"]
    motif = brief.get("visual_motif", "")

    if style in ("editorial", "minimal"):
        _light_geometric_background(slide, theme)
    else:
        _geometric_background(slide, theme)
    _apply_motif(slide, theme, motif)

    _add_text_box(slide, Inches(0), Inches(2.0), SLIDE_W, Inches(1.2),
                  "Thank You", 52, theme.white, bold=True, align=PP_ALIGN.CENTER, font=title_font)

    _add_text_box(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(0.6),
                  "AI-Powered Analysis  •  Data-Driven Decisions", 16, theme.muted,
                  align=PP_ALIGN.CENTER, italic=True, font=body_font)

    # Centered accent line
    _add_rect(slide, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.03), theme.accent(0))

    _add_text_box(slide, Inches(0), Inches(5.0), SLIDE_W, Inches(0.5),
                  f"Analysis powered by NexusMind Studio  •  Theme: {brief.get('theme', 'Professional')}",
                  10, theme.dark_muted, align=PP_ALIGN.CENTER, font=body_font)

    _bottom_strip(slide, theme, alpha=0.45 if style in ("editorial", "minimal") else 0.85)


def _build_slide_plan(
    narrative: str,
    sections: List[Dict[str, Any]],
    insights: List[Dict[str, Any]],
    design_spec: Dict[str, Any],
    chart_paths: List[str],
) -> Dict[str, Any]:
    """Decide deck length from content depth instead of a fixed template list."""
    density = design_spec.get("slide_density", "medium")
    significant = [
        ins for ins in insights
        if _safe_score(ins.get("significance_score")) >= 0.35 or ins.get("data_evidence")
    ]
    if not significant:
        significant = insights

    density_caps = {"minimal": 3, "medium": 5, "rich": 8}
    insight_cap = density_caps.get(density, 5)
    insight_count = min(len(significant), insight_cap)

    valid_chart_count = sum(1 for p in chart_paths if _resolve_png_path(p))
    highlighted_charts = sum(
        1 for spec in design_spec.get("chart_specs", [])
        if isinstance(spec, dict) and spec.get("highlight_insight")
    )
    if density == "minimal":
        chart_count = min(valid_chart_count, max(1, highlighted_charts or 1)) if valid_chart_count else 0
    elif density == "rich":
        chart_count = min(valid_chart_count, max(2, insight_count, highlighted_charts))
    else:
        chart_count = min(valid_chart_count, max(1, min(insight_count, 4), highlighted_charts))

    has_recommendations = bool(_extract_recommendation_bullets(sections))
    has_story = bool(design_spec.get("storytelling_arc")) and (insight_count >= 2 or chart_count >= 2)
    include_summary = bool(narrative.strip() or sections or insights)
    include_closing = insight_count + chart_count >= 3 or has_recommendations

    return {
        "include_summary": include_summary,
        "include_story": has_story,
        "insight_count": insight_count,
        "chart_count": chart_count,
        "include_recommendations": has_recommendations,
        "include_closing": include_closing,
    }


def _safe_score(value: Any) -> float:
    try:
        return float(value or 0)
    except (TypeError, ValueError):
        return 0.0


def _resolve_subtitle(narrative: str, sections: List[Dict[str, Any]], design_spec: Dict[str, Any]) -> str:
    configured = design_spec.get("subtitle") or design_spec.get("deck_subtitle")
    if configured:
        return str(configured)
    if narrative:
        first_sentence = narrative.strip().split(". ")[0].strip()
        if first_sentence:
            return first_sentence if first_sentence.endswith(".") else f"{first_sentence}."
    if sections:
        content = str(sections[0].get("content", "")).strip()
        if content:
            first_sentence = content.split(". ")[0].strip()
            return first_sentence[:150]
    return "A data-backed narrative of patterns, evidence, and recommended action."


def _extract_recommendation_bullets(sections: List[Dict[str, Any]]) -> List[str]:
    rec_section = next((s for s in sections if "recommend" in s.get("title", "").lower()), None)
    if rec_section:
        return _extract_bullets(rec_section.get("content", ""))
    if sections:
        return _extract_bullets(sections[-1].get("content", ""))
    return []


# ═══════════════════════════════════════════════════════════
# Orchestrator
# ═══════════════════════════════════════════════════════════

def generate_slides(
    title: str,
    narrative: str,
    sections: List[Dict[str, Any]],
    insights: List[Dict[str, Any]],
    design_spec: Dict[str, Any],
    chart_paths: List[str],
    job_id: str,
) -> str:
    """Generate a professional PPTX from analysis results."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    palette = design_spec.get("color_palette", ["#4f81bd", "#9cbb58", "#f79646", "#8064a2"])
    theme_name = design_spec.get("theme", "generic")
    style_profile = _resolve_template(design_spec, theme_name)
    style_name = style_profile["name"]
    theme = ThemeEngine(palette, theme_name, mode=style_profile.get("mode", "dark"))
    typography = _resolve_typography(design_spec, style_name)
    subtitle = _resolve_subtitle(narrative, sections, design_spec)

    motif = design_spec.get("visual_motif", "")
    plan = _build_slide_plan(narrative, sections, insights, design_spec, chart_paths)

    slide_steps = [("title", lambda: _title_slide(prs, title, subtitle, theme, palette, design_spec, style_name, typography))]
    if plan["include_summary"]:
        slide_steps.append(("summary", lambda: _executive_summary_slide(prs, narrative, theme, sections, style_name, typography)))
    if plan["include_story"]:
        slide_steps.append(("story_arc", lambda: _story_arc_slide(prs, design_spec, theme, style_name, typography)))
    if insights and plan["insight_count"]:
        slide_steps.append((
            "insights",
            lambda: _insight_slides(
                prs, insights, chart_paths, theme, palette, design_spec, style_name, typography,
                insight_limit=plan["insight_count"],
            ),
        ))
    if chart_paths and plan["chart_count"]:
        slide_steps.append((
            "charts",
            lambda: _chart_slides(
                prs, chart_paths, theme, palette, style_name, typography, motif,
                chart_specs=design_spec.get("chart_specs", []), insights=insights,
                chart_limit=plan["chart_count"],
            ),
        ))
    if plan["include_recommendations"]:
        slide_steps.append(("recommendations", lambda: _recommendations_slide(prs, sections, theme, palette, style_name, typography, motif)))
    if plan["include_closing"]:
        slide_steps.append(("closing", lambda: _closing_slide(prs, theme, design_spec, style_name, typography)))

    for element, builder in slide_steps:
        try:
            builder()
        except Exception as exc:
            logger.log_error(f"Slide generation failed for: {element}", exc, job_id=job_id)

    # Persist
    output_dir = os.path.join("data", "slides")
    os.makedirs(output_dir, exist_ok=True)
    path = os.path.join(output_dir, f"{job_id}.pptx")
    prs.save(path)
    logger.log_operation("Slides generated", job_id=job_id, slide_count=len(prs.slides), path=path, slide_plan=plan)
    return path


# ═══════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════

def _hex(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _to_hex(c: RGBColor) -> str:
    return f"{c[0]:02x}{c[1]:02x}{c[2]:02x}"


def _opacity(c: RGBColor, alpha: float) -> RGBColor:
    """Return a color that visually approximates the base color at alpha
    when rendered over a dark background. Used for surfaces/cards."""
    return RGBColor(
        max(0, min(255, int(c[0] * alpha + 26 * (1 - alpha)))),
        max(0, min(255, int(c[1] * alpha + 26 * (1 - alpha)))),
        max(0, min(255, int(c[2] * alpha + 30 * (1 - alpha)))),
    )


def _set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _section_label(slide, text: str, theme: ThemeEngine, font: str):
    """Small uppercase section label at top of slide."""
    _add_text_box(slide, MARGIN, Inches(0.25), Inches(6), Inches(0.3),
                  text, 9, theme.accent(0), bold=True, font=font)


def _resolve_png_path(path: str) -> str | None:
    """Get the PNG version of a chart path if available."""
    if not path:
        return None
    if path.endswith(".png"):
        return path if os.path.exists(path) else None
    if path.endswith(".json"):
        png = path.replace(".json", ".png")
        return png if os.path.exists(png) else None
    return path if os.path.exists(path) else None


def _clean_markdown_text(text: Any) -> str:
    """Remove lightweight Markdown artifacts and boilerplates before placing text in PPTX boxes."""
    cleaned = str(text or "").replace("\r", " ").replace("\n", " ")
    
    # Strip common LLM placeholders
    placeholders = [r"\[Insert text here\]", r"\[Insert.*?\]", r"\bTBD\b", r"Feature block", r"\[.*?\]"]
    for p in placeholders:
        cleaned = re.sub(p, "", cleaned, flags=re.IGNORECASE)
        
    cleaned = re.sub(r"\*\*(.*?)\*\*", r"\1", cleaned)
    cleaned = re.sub(r"__(.*?)__", r"\1", cleaned)
    cleaned = re.sub(r"`([^`]*)`", r"\1", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned)
    return cleaned.strip(" -*•▸")


def _extract_bullets(md_text: str) -> List[str]:
    """Extract bullet points from markdown text, dropping empty or filler points."""
    bullets: List[str] = []
    for line in md_text.split("\n"):
        stripped = line.strip()
        bullet_text = ""
        if stripped.startswith(("- ", "* ", "• ", "▸ ")):
            bullet_text = _clean_markdown_text(stripped[2:].strip())
        elif stripped.startswith(tuple("0123456789")) and ". " in stripped[:4]:
            bullet_text = _clean_markdown_text(stripped.split(". ", 1)[1].strip())
            
        if bullet_text and len(bullet_text) > 3: # Skip empty or tiny filler
            bullets.append(bullet_text)
            
    return bullets[:8]
